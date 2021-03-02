#!/usr/bin/env python
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt

pr1 = Presentation()

slide1_register = pr1.slide_layouts[6]

slide1 = pr1.slides.add_slide(slide1_register)

left = top = width = height = Inches(2)
txBox = slide1.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

#print(b'\xf0\xaa\x9c\xb6')
#msg = b'\xe7\x9a\x84'
msg = b'\xf0\xaa\x9c\xb6'
tf.text = msg.decode()
#tf.text = "十"

#title1 = slide1.shapes.title
#title1.text = "\u7684\u6511"

#body_shape = shapes.placeholders[1]


pr1.save('PythonPPT.pptx')

